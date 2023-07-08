import pptx
import datetime
from shutil import copyfile
import re

from .settings import powerpoint_vorlage

from home.models import WeekMotto

def get_event_date_from_agenda(agenda_dictionary):
    return agenda_dictionary["data"]["name"][0:10]


def create_pp_file(event_name, event_folder):
    new_filename = event_folder + event_name + '.pptx'
    copyfile(powerpoint_vorlage, new_filename)
    return new_filename


def get_weekvers_from_agenda(agenda_dictionary):
    event_date = get_event_date_from_agenda(agenda_dictionary)
    datetime_object = datetime.datetime.strptime(event_date, '%d.%m.%Y')
    weekmottos = WeekMotto.objects.filter(date=datetime_object)
    vers = ''
    if len(weekmottos) == 1:
        for weekmotto in weekmottos:
            return weekmotto.motto_luther_modern

    for weekmotto in weekmottos:
        if weekmotto.motto_luther_modern == vers:
            return weekmotto.motto_luther_modern
        else:
            vers = weekmotto.motto_luther_modern

    for weekmotto in weekmottos:
        if 'Sonntag' in weekmotto.title:
            return weekmotto.motto_luther_modern
    
    for weekmotto in weekmottos:
        return weekmotto.motto_luther_modern
    
    return 'kein Vers für dieses Datum gefunden'


def set_informations_in_pp_by_placeholder(pp_filename, informations):
    if not informations is None:
        if informations != '':
            if "\n-" in informations:
                informations = re.sub(r'^- ', '', informations)
                informations_list = informations.split('\n-')
            else:
                informations = re.sub(r'^- ', '', informations)
                informations_list = [informations]
        else:
            informations_list = ["Noch keine Informationen vorhanden!!"]
    else:
        informations_list = ["Noch keine Informationen vorhanden!!"]

    prs = pptx.Presentation(pp_filename)

    amount_of_information_on_a_slide = 3
    for information in informations_list:
        index = informations_list.index(information)
        layout_number = (index//amount_of_information_on_a_slide) + 1
        layout_name = 'Informationen_' + str(layout_number)
        information = information.lstrip()
        for slide in prs.slides:
            if layout_name == slide.slide_layout.name:
                placeholders = slide.shapes.placeholders
                placeholder = placeholders[13]
                placeholder.text = placeholder.text + information + '\n'

    xml_slides = prs.slides._sldIdLst  
    slides = list(xml_slides)
    counter = 0

    for slide in prs.slides:
        if 'Informationen_' in slide.slide_layout.name:
            placeholders = slide.shapes.placeholders
            placeholder = placeholders[13]
            if placeholder.text == '':
                index = prs.slides.index(slide) + counter
                xml_slides.remove(slides[index])
                counter += 1

    prs.save(pp_filename)


def set_weekvers_in_pp_by_placeholder(pp_filename, agenda_dictionary):
    weekvers = get_weekvers_from_agenda(agenda_dictionary)

    prs = pptx.Presentation(pp_filename)
    for slide in prs.slides:
        if slide.slide_layout.name == "Wochenspruch":
            placeholders = slide.shapes.placeholders
            placeholder = placeholders[13]
            placeholder.text = weekvers

    prs.save(pp_filename)